"""文件内容解析模块 — 支持 PDF/PPTX/DOCX/XLSX/图片/纯文本"""
import base64
import io
import os
from typing import Optional

# 支持的文件类型白名单
ALLOWED_EXTENSIONS = {
    # 文档
    '.pdf', '.pptx', '.docx', '.xlsx',
    # 图片
    '.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp',
    # 纯文本
    '.txt', '.md', '.py', '.js', '.ts', '.jsx', '.tsx', '.vue',
    '.json', '.csv', '.html', '.css', '.xml', '.yaml', '.yml',
    '.log', '.sh', '.bat', '.java', '.c', '.cpp', '.h', '.cs',
    '.sql', '.r', '.go', '.rs', '.swift', '.kt', '.toml', '.ini', '.cfg',
}

TEXT_EXTENSIONS = {
    '.txt', '.md', '.py', '.js', '.ts', '.jsx', '.tsx', '.vue',
    '.json', '.csv', '.html', '.css', '.xml', '.yaml', '.yml',
    '.log', '.sh', '.bat', '.java', '.c', '.cpp', '.h', '.cs',
    '.sql', '.r', '.go', '.rs', '.swift', '.kt', '.toml', '.ini', '.cfg',
}

IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp'}

MAX_FILE_SIZE = 20 * 1024 * 1024  # 20MB

# ---- OCR 引擎（双引擎 fallback：tesseract → easyocr） ----
_easyocr_reader = None
_tesseract_available = None
_TESSERACT_CMD = None  # None 表示让 pytesseract 自动检测
_TESSDATA_DIR = None  # 自定义 tessdata 目录（用于存放额外下载的语言包）

# 常见平台上的 Tesseract 安装路径（按优先级排列）
_TESSERACT_CANDIDATES = [
    # 自动检测（系统 PATH）
    'tesseract',
    # Windows 默认路径
    r'C:/Program Files/Tesseract-OCR/tesseract.exe',
    r'C:/Program Files (x86)/Tesseract-OCR/tesseract.exe',
    # Linux 默认路径
    '/usr/bin/tesseract',
    '/usr/local/bin/tesseract',
    # macOS Homebrew
    '/opt/homebrew/bin/tesseract',
    '/usr/local/bin/tesseract',
]


def _setup_tesseract_env():
    """配置 Tesseract 可执行文件路径和 tessdata 目录"""
    import os as _os
    # 如果用户目录下有 tessdata（含额外语言包），优先使用
    user_tessdata = _os.path.expanduser('~/tessdata')
    if _os.path.isdir(user_tessdata):
        _os.environ.setdefault('TESSDATA_PREFIX', user_tessdata)


def _check_tesseract() -> bool:
    """检测 Tesseract 是否可用（离线 OCR 引擎），跨平台自动查找安装路径"""
    global _tesseract_available, _TESSERACT_CMD
    if _tesseract_available is None:
        try:
            import pytesseract
            import shutil as _shutil
            _setup_tesseract_env()
            # 如果尚未确定路径，逐个尝试候选路径
            if _TESSERACT_CMD is None:
                for _candidate in _TESSERACT_CANDIDATES:
                    if _shutil.which(_candidate):
                        _TESSERACT_CMD = _candidate
                        break
            if _TESSERACT_CMD:
                pytesseract.pytesseract.tesseract_cmd = _TESSERACT_CMD
            pytesseract.get_tesseract_version()
            _tesseract_available = True
        except Exception:
            _tesseract_available = False
    return _tesseract_available


def _get_tesseract_lang() -> str:
    """返回可用的 Tesseract 语言设置，优先中英混合，回退纯英文"""
    try:
        import pytesseract
        _setup_tesseract_env()
        pytesseract.pytesseract.tesseract_cmd = _TESSERACT_CMD
        langs = pytesseract.get_languages()
        if 'chi_sim' in langs:
            return 'chi_sim+eng'
        return 'eng'
    except Exception:
        return 'eng'


def _ocr_tesseract(file_bytes: bytes) -> str:
    """使用 pytesseract 离线 OCR"""
    if not _check_tesseract():
        return ""
    try:
        import pytesseract
        from PIL import Image
        _setup_tesseract_env()
        pytesseract.pytesseract.tesseract_cmd = _TESSERACT_CMD
        image = Image.open(io.BytesIO(file_bytes))
        if image.mode in ('RGBA', 'P', 'LA'):
            image = image.convert('RGB')
        lang = _get_tesseract_lang()
        text = pytesseract.image_to_string(image, lang=lang)
        return text.strip()
    except Exception:
        return ""


def _get_easyocr_reader():
    """懒加载 easyocr Reader --- 任何初始化失败都返回 None，不抛异常"""
    global _easyocr_reader
    if _easyocr_reader is None:
        try:
            import easyocr
            _easyocr_reader = easyocr.Reader(['ch_sim', 'en'], gpu=False, verbose=False)
        except Exception:
            # ImportError / 模型下载网络错误 / 任何初始化失败 → 禁用 easyocr
            _easyocr_reader = None
            return None
    return _easyocr_reader


def _ocr_easyocr(file_bytes: bytes) -> str:
    """使用 easyocr 在线 OCR（首次运行需下载模型）"""
    reader = _get_easyocr_reader()
    if reader is None:
        return ""
    try:
        from PIL import Image
        image = Image.open(io.BytesIO(file_bytes))
        if image.mode in ('RGBA', 'P', 'LA'):
            image = image.convert('RGB')
        results = reader.readtext(image)
        text_lines = [item[1] for item in results if item[2] > 0.2]
        return '\n'.join(text_lines)
    except Exception:
        return ""


def _ocr_image(file_bytes: bytes) -> str:
    """对图片二进制数据进行 OCR 识别，依次尝试 tesseract → easyocr"""
    # 优先 Tesseract（离线，速度快）
    text = _ocr_tesseract(file_bytes)
    if text:
        return text
    # 回退 easyocr（在线，首次需下载模型）
    text = _ocr_easyocr(file_bytes)
    if text:
        return text
    return ""


def get_file_type(filename: str) -> str:
    """根据扩展名返回文件类型描述"""
    ext = os.path.splitext(filename)[1].lower()
    type_map = {
        '.pdf': 'PDF文档', '.pptx': 'PPT演示文稿', '.docx': 'Word文档',
        '.xlsx': 'Excel表格', '.png': '图片', '.jpg': '图片', '.jpeg': '图片',
        '.gif': '图片', '.bmp': '图片', '.webp': '图片',
        '.py': 'Python代码', '.js': 'JavaScript代码', '.ts': 'TypeScript代码',
        '.java': 'Java代码', '.c': 'C代码', '.cpp': 'C++代码',
        '.sql': 'SQL脚本', '.html': 'HTML文件', '.css': 'CSS文件',
        '.json': 'JSON数据', '.csv': 'CSV数据', '.md': 'Markdown文档',
    }
    return type_map.get(ext, f'{ext}文件')


def parse_file(file_bytes: bytes, filename: str) -> dict:
    """
    解析上传文件，提取文本内容。
    返回 {"text": str, "file_type": str, "base64": Optional[str]}
    """
    ext = os.path.splitext(filename)[1].lower()

    if ext not in ALLOWED_EXTENSIONS:
        raise ValueError(f"不支持的文件格式: {ext}。支持: {', '.join(sorted(ALLOWED_EXTENSIONS))}")

    if len(file_bytes) > MAX_FILE_SIZE:
        raise ValueError(f"文件大小超过限制（最大 20MB）")

    file_type = get_file_type(filename)
    base64_data = None

    # --- 图片：OCR 提取文字 + base64 编码 ---
    if ext in IMAGE_EXTENSIONS:
        base64_data = base64.b64encode(file_bytes).decode('utf-8')
        mime_map = {'.png': 'image/png', '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg',
                    '.gif': 'image/gif', '.bmp': 'image/bmp', '.webp': 'image/webp'}
        mime = mime_map.get(ext, 'image/png')

        # OCR 识别文字（tesseract → easyocr 双引擎 fallback）
        ocr_text = _ocr_image(file_bytes)
        if ocr_text:
            text = f'[图片文件: {filename}]\n（以下为图片OCR识别内容，请基于此内容回答用户问题）\n\n{ocr_text}'
        else:
            # 两个 OCR 引擎都不可用时，提示用户
            text = (
                f'[图片文件: {filename}]\n'
                f'⚠️ 图片文字识别失败。请尝试以下方案：\n'
                f'1. 安装 Tesseract OCR（推荐，离线）：下载地址 https://github.com/UB-Mannheim/tesseract/wiki\n'
                f'   安装时勾选中文语言包，然后重启后端服务即可自动识别图片\n'
                f'2. 或将图片内容以文字形式直接输入'
            )

        return {
            'text': text,
            'file_type': file_type,
            'base64': f'data:{mime};base64,{base64_data}',
        }

    # --- PDF ---
    if ext == '.pdf':
        text = _parse_pdf(file_bytes)
        return {'text': text, 'file_type': file_type, 'base64': None}

    # --- PPTX ---
    if ext == '.pptx':
        text = _parse_pptx(file_bytes)
        return {'text': text, 'file_type': file_type, 'base64': None}

    # --- DOCX ---
    if ext == '.docx':
        text = _parse_docx(file_bytes)
        return {'text': text, 'file_type': file_type, 'base64': None}

    # --- XLSX ---
    if ext == '.xlsx':
        text = _parse_xlsx(file_bytes)
        return {'text': text, 'file_type': file_type, 'base64': None}

    # --- 纯文本 ---
    if ext in TEXT_EXTENSIONS:
        text = _parse_text(file_bytes, filename)
        return {'text': text, 'file_type': file_type, 'base64': None}

    raise ValueError(f"未处理的文件格式: {ext}")


def _parse_pdf(data: bytes) -> str:
    """使用 PyMuPDF 提取 PDF 文本"""
    import fitz  # PyMuPDF
    doc = fitz.open(stream=data, filetype='pdf')
    parts = []
    for page_num, page in enumerate(doc, 1):
        text = page.get_text('text')
        if text.strip():
            parts.append(f'--- 第{page_num}页 ---\n{text.strip()}')
    doc.close()
    return '\n\n'.join(parts) if parts else '(PDF内容为空)'


def _parse_pptx(data: bytes) -> str:
    """使用 python-pptx 提取演示文稿文本"""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    t = paragraph.text.strip()
                    if t:
                        slide_texts.append(t)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        slide_texts.append(row_text)
        if slide_texts:
            parts.append(f'--- 幻灯片{slide_num} ---\n' + '\n'.join(slide_texts))
    return '\n\n'.join(parts) if parts else '(PPT内容为空)'


def _parse_docx(data: bytes) -> str:
    """使用 python-docx 提取 Word 文档文本"""
    from docx import Document
    doc = Document(io.BytesIO(data))
    parts = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)
    # 提取表格
    for table in doc.tables:
        for row in table.rows:
            row_text = ' | '.join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)
    return '\n'.join(parts) if parts else '(文档内容为空)'


def _parse_xlsx(data: bytes) -> str:
    """使用 openpyxl 提取 Excel 表格内容"""
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f'--- Sheet: {sheet_name} ---')
        for row in ws.iter_rows(values_only=True):
            row_text = ' | '.join(str(c) if c is not None else '' for c in row)
            if row_text.strip():
                parts.append(row_text)
    wb.close()
    return '\n'.join(parts) if parts else '(表格内容为空)'


def _parse_text(data: bytes, filename: str) -> str:
    """读取纯文本文件"""
    for encoding in ['utf-8', 'gbk', 'gb2312', 'latin-1']:
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode('utf-8', errors='replace')
